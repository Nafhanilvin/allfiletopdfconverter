from flask import Flask, request, send_file, jsonify, render_template
from flask_cors import CORS
import os
from werkzeug.utils import secure_filename
import io
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from reportlab.lib.pagesizes import letter, A4, landscape
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from PIL import Image
import tempfile
import os

app = Flask(__name__)
CORS(app)

# Configuration
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
ALLOWED_EXTENSIONS = {'xlsx', 'xls', 'xlsm', 'pptx', 'ppt', 'docx', 'doc', 'txt'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_file_type(filename):
    """Determine if file is Excel, PowerPoint, Word, or Text"""
    ext = filename.rsplit('.', 1)[1].lower()
    if ext in {'xlsx', 'xls', 'xlsm'}:
        return 'excel'
    elif ext in {'pptx', 'ppt'}:
        return 'powerpoint'
    elif ext in {'docx', 'doc'}:
        return 'word'
    elif ext == 'txt':
        return 'text'
    return None

def excel_to_pdf(excel_file):
    """Convert Excel file to PDF"""
    # Load the workbook
    wb = load_workbook(excel_file, data_only=True)
    
    # Determine orientation based on data width
    # First, check how many columns we have
    max_cols_in_any_sheet = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            max_cols_in_any_sheet = max(max_cols_in_any_sheet, len([c for c in row if c is not None]))
    
    # Use landscape for sheets with more than 8 columns, portrait otherwise
    use_landscape = max_cols_in_any_sheet > 8
    pagesize = landscape(A4) if use_landscape else A4
    
    # Create PDF in memory with appropriate orientation
    pdf_buffer = io.BytesIO()
    doc = SimpleDocTemplate(pdf_buffer, pagesize=pagesize, leftMargin=0.25*inch, rightMargin=0.25*inch, topMargin=0.25*inch, bottomMargin=0.25*inch)
    
    elements = []
    styles = getSampleStyleSheet()
    
    # Process each sheet
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        
        # Add sheet name as header
        if len(wb.sheetnames) > 1:
            header = Paragraph(f"<b>{sheet_name}</b>", styles['Heading2'])
            elements.append(header)
            elements.append(Spacer(1, 0.1*inch))
        
        # Get data from sheet
        data = []
        for row in ws.iter_rows(values_only=True):
            # Convert None to empty string and all values to strings
            row_data = [str(cell) if cell is not None else '' for cell in row]
            data.append(row_data)
        
        if data:
            # Remove completely empty rows at the end
            while data and all(cell == '' for cell in data[-1]):
                data.pop()
            
            # Remove completely empty columns at the end
            if data:
                max_cols = max(len(row) for row in data)
                for row in data:
                    while len(row) < max_cols:
                        row.append('')
                
                # Find and remove completely empty trailing columns
                while max_cols > 0:
                    if all(row[max_cols-1] == '' for row in data):
                        for row in data:
                            if len(row) > max_cols - 1:
                                row.pop()
                        max_cols -= 1
                    else:
                        break
            
            if data and any(any(cell != '' for cell in row) for row in data):
                max_cols = max(len(row) for row in data) if data else 0
                
                if max_cols > 0:
                    # Calculate optimal column widths based on page orientation
                    if use_landscape:
                        available_width = landscape(A4)[0] - 0.5*inch
                    else:
                        available_width = A4[0] - 0.5*inch
                    
                    col_width = available_width / max_cols
                    
                    # Limit minimum column width to ensure readability
                    min_col_width = 0.35*inch
                    if col_width < min_col_width:
                        col_width = min_col_width
                    
                    col_widths = [col_width] * max_cols
                    
                    # Create table
                    table = Table(data, colWidths=col_widths)
                    
                    # Style the table
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4472C4')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 7),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 5),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#D9E2F3')),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#A6B4D0')),
                        ('FONTSIZE', (0, 1), (-1, -1), 5),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#E7EEF7')]),
                    ]))
                    
                    elements.append(table)
                    elements.append(Spacer(1, 0.15*inch))
    
    # Build PDF
    if not elements:
        elements.append(Paragraph("No data found in the spreadsheet.", styles['Normal']))
    
    doc.build(elements)
    pdf_buffer.seek(0)
    
    return pdf_buffer

def pptx_to_pdf(pptx_file):
    """Convert PowerPoint file to PDF with images"""
    try:
        # Load the presentation
        prs = Presentation(pptx_file)
        
        # Create PDF in memory
        pdf_buffer = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
        
        elements = []
        styles = getSampleStyleSheet()
        temp_files = []  # Keep track of all temp files
        
        # Get page dimensions
        page_width = A4[0] - 2*inch  # Leave 1 inch margin on each side
        
        # Process each slide
        for slide_idx, slide in enumerate(prs.slides, 1):
            try:
                # Add slide number header
                header = Paragraph(f"<b>Slide {slide_idx}</b>", styles['Heading2'])
                elements.append(header)
                elements.append(Spacer(1, 0.15*inch))
                
                # Extract and add images from slide
                slide_has_images = False
                
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "image"):
                            # Get the image
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Save to temporary file
                            with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{image.ext}') as tmp_file:
                                tmp_file.write(image_bytes)
                                tmp_path = tmp_file.name
                            
                            temp_files.append(tmp_path)  # Track temp file
                            
                            try:
                                # Open with PIL to get dimensions and ensure it's valid
                                pil_img = Image.open(tmp_path)
                                img_width, img_height = pil_img.size
                                pil_img.close()
                                
                                # Calculate scaling to fit page width while maintaining aspect ratio
                                aspect = img_height / float(img_width)
                                display_width = min(page_width, 5*inch)  # Max 5 inches wide
                                display_height = display_width * aspect
                                
                                # Add image to PDF - use the temp path directly
                                img = RLImage(tmp_path, width=display_width, height=display_height)
                                elements.append(img)
                                elements.append(Spacer(1, 0.1*inch))
                                slide_has_images = True
                            except Exception as e:
                                print(f"Error processing image in slide {slide_idx}: {e}")
                                continue
                    
                    except Exception as e:
                        print(f"Error extracting image from shape: {e}")
                        continue
                
                # Extract text from shapes
                slide_text = []
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    except Exception as e:
                        print(f"Error extracting text from shape: {e}")
                        continue
                
                # Add text content
                if slide_text:
                    for text in slide_text:
                        try:
                            # Clean and format text
                            text = text.replace('\r\n', '<br/>').replace('\n', '<br/>')
                            para = Paragraph(text, styles['Normal'])
                            elements.append(para)
                            elements.append(Spacer(1, 0.1*inch))
                        except Exception as e:
                            print(f"Error adding text to PDF: {e}")
                            continue
                elif not slide_has_images:
                    # If no text and no images, add placeholder
                    para = Paragraph("<i>[Empty slide]</i>", styles['Normal'])
                    elements.append(para)
                    elements.append(Spacer(1, 0.1*inch))
                
                # Add page break between slides (except last one)
                if slide_idx < len(prs.slides):
                    elements.append(PageBreak())
            
            except Exception as e:
                print(f"Error processing slide {slide_idx}: {e}")
                continue
        
        # Build PDF
        if not elements:
            # If no elements were added, add a message
            elements.append(Paragraph("No content could be extracted from the presentation.", styles['Normal']))
        
        try:
            doc.build(elements)
        except Exception as e:
            print(f"Error building PDF: {e}")
            raise
        
        pdf_buffer.seek(0)
        
        # Clean up temp files after PDF is built
        for tmp_path in temp_files:
            try:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
            except Exception as e:
                print(f"Error deleting temp file {tmp_path}: {e}")
        
        return pdf_buffer
    
    except Exception as e:
        print(f"Fatal error in pptx_to_pdf: {e}")
        raise

def docx_to_pdf(docx_file):
    """Convert Word document to PDF"""
    try:
        # Load the document
        doc = Document(docx_file)
        
        # Create PDF in memory
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=A4, leftMargin=0.75*inch, rightMargin=0.75*inch, topMargin=0.75*inch, bottomMargin=0.75*inch)
        
        elements = []
        styles = getSampleStyleSheet()
        
        # Get styles
        title_style = styles['Heading1']
        heading_style = styles['Heading2']
        body_style = styles['Normal']
        
        # Process paragraphs
        for para in doc.paragraphs:
            try:
                if not para.text.strip():
                    # Add space for empty paragraphs
                    elements.append(Spacer(1, 0.1*inch))
                else:
                    # Detect paragraph style/formatting
                    text = para.text.strip()
                    
                    # Get paragraph style level to determine formatting
                    style_name = para.style.name if para.style else 'Normal'
                    
                    # Create styled paragraph
                    if 'Heading 1' in style_name or 'Title' in style_name:
                        para_obj = Paragraph(f"<b><font size=14>{text}</font></b>", title_style)
                    elif 'Heading 2' in style_name or 'Heading' in style_name:
                        para_obj = Paragraph(f"<b><font size=12>{text}</font></b>", heading_style)
                    else:
                        # Apply formatting from the paragraph
                        formatted_text = text
                        if para.runs:
                            formatted_parts = []
                            for run in para.runs:
                                run_text = run.text
                                if run.bold:
                                    run_text = f"<b>{run_text}</b>"
                                if run.italic:
                                    run_text = f"<i>{run_text}</i>"
                                if run.underline:
                                    run_text = f"<u>{run_text}</u>"
                                formatted_parts.append(run_text)
                            formatted_text = ''.join(formatted_parts)
                        
                        para_obj = Paragraph(formatted_text, body_style)
                    
                    elements.append(para_obj)
                    elements.append(Spacer(1, 0.05*inch))
            
            except Exception as e:
                print(f"Error processing paragraph: {e}")
                continue
        
        # Process tables
        for table in doc.tables:
            try:
                # Extract table data
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        row_data.append(cell_text)
                    table_data.append(row_data)
                
                if table_data:
                    # Create table
                    pdf_table = Table(table_data)
                    
                    # Style the table
                    pdf_table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4472C4')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 9),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#D9E2F3')),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#A6B4D0')),
                        ('FONTSIZE', (0, 1), (-1, -1), 8),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#E7EEF7')]),
                    ]))
                    
                    elements.append(pdf_table)
                    elements.append(Spacer(1, 0.2*inch))
            
            except Exception as e:
                print(f"Error processing table: {e}")
                continue
        
        # Build PDF
        if not elements:
            elements.append(Paragraph("No content found in the document.", styles['Normal']))
        
        pdf_doc.build(elements)
        pdf_buffer.seek(0)
        
        return pdf_buffer
    
    except Exception as e:
        print(f"Fatal error in docx_to_pdf: {e}")
        raise

def txt_to_pdf(txt_file):
    """Convert text file to PDF"""
    try:
        # Read the text file
        content = txt_file.read().decode('utf-8', errors='replace')
        
        # Create PDF in memory
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=A4, leftMargin=0.75*inch, rightMargin=0.75*inch, topMargin=0.75*inch, bottomMargin=0.75*inch)
        
        elements = []
        styles = getSampleStyleSheet()
        body_style = styles['Normal']
        
        # Split content into lines and process
        lines = content.split('\n')
        
        for line in lines:
            try:
                line = line.rstrip()
                
                if not line.strip():
                    # Add space for empty lines
                    elements.append(Spacer(1, 0.1*inch))
                else:
                    # Escape special characters for PDF
                    safe_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    para = Paragraph(safe_line, body_style)
                    elements.append(para)
                    elements.append(Spacer(1, 0.05*inch))
            
            except Exception as e:
                print(f"Error processing line: {e}")
                continue
        
        # Build PDF
        if not elements:
            elements.append(Paragraph("No content found in the text file.", styles['Normal']))
        
        pdf_doc.build(elements)
        pdf_buffer.seek(0)
        
        return pdf_buffer
    
    except Exception as e:
        print(f"Fatal error in txt_to_pdf: {e}")
        raise

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert():
    try:
        # Check if file was uploaded
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': 'Invalid file type. Please upload an Excel or PowerPoint file (.xlsx, .xls, .xlsm, .pptx, .ppt)'}), 400
        
        # Determine file type and convert accordingly
        file_type = get_file_type(file.filename)
        
        if file_type == 'excel':
            pdf_buffer = excel_to_pdf(file)
        elif file_type == 'powerpoint':
            pdf_buffer = pptx_to_pdf(file)
        elif file_type == 'word':
            pdf_buffer = docx_to_pdf(file)
        elif file_type == 'text':
            pdf_buffer = txt_to_pdf(file)
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
        
        # Generate output filename
        original_filename = secure_filename(file.filename)
        pdf_filename = os.path.splitext(original_filename)[0] + '.pdf'
        
        return send_file(
            pdf_buffer,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=pdf_filename
        )
    
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"Conversion error: {error_trace}")
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy'}), 200

if __name__ == '__main__':
    import os
    port = int(os.environ.get('PORT', 5000))
    app.run(debug=False, host='0.0.0.0', port=port)
